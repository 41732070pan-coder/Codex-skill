# -*- coding: utf-8 -*-
"""Regenerate the skill_intro_deck sample PPTX from content_outline.json."""

from __future__ import annotations

import hashlib
import json
import random
from collections import Counter
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from reportlab.graphics import renderPM
from svglib.svglib import svg2rlg
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

EXAMPLE_ROOT = Path(__file__).resolve().parent
SKILL_ROOT = EXAMPLE_ROOT.parents[1]
ASSET_ROOT = SKILL_ROOT / "assets" / "chinese_traditional_color_style"
OUT_DIR = EXAMPLE_ROOT / "generated"
OUTLINE_PATH = EXAMPLE_ROOT / "content_outline.json"

W, H = 13.333, 7.5

COLORS = {
    "paper": "FFFBF0",
    "white": "FFFFFF",
    "soft": "F2ECDE",
    "moon": "D6ECF0",
    "mist": "E9F1F6",
    "indigo": "177CB0",
    "indigo_deep": "065279",
    "bluegray": "A1AFC9",
    "cang": "75878A",
    "ink": "161823",
    "ink_soft": "50616D",
    "silver": "BACAC6",
    "bamboo": "789262",
    "jade": "2EDFA3",
    "gold": "EACD76",
    "apricot": "FFA631",
    "rouge": "9D2933",
    "rouge_soft": "DB5A6B",
    "purple": "8D4BBB",
    "587558": "587558",
    "FDD000": "FDD000",
    "AD986E": "AD986E",
    "6F8B6C": "6F8B6C",
    "496B61": "496B61",
    "364A88": "364A88",
    "6A4D31": "6A4D31",
    "443459": "443459",
    "BD0E2B": "BD0E2B",
}

FONTS = {"body": "Microsoft YaHei", "title": "SimSun", "mono": "Consolas"}

# svglib + renderPM rasterize onto an opaque canvas (default white). Converting to RGBA
# alone keeps that canvas visible; we must knock out the detected background color.
RASTER_CACHE_VERSION = "transparent_v2"


def strip_solid_background(img: Image.Image, tolerance: int = 28) -> Image.Image:
    """Turn the raster canvas color (usually white) into real alpha."""
    img = img.convert("RGBA")
    w, h = img.size
    if w < 2 or h < 2:
        return img

    px = img.load()
    corner_colors: list[tuple[int, int, int]] = []

    def sample_corner(cx: int, cy: int, radius: int = 6) -> None:
        for y in range(max(0, cy - radius), min(h, cy + radius + 1)):
            for x in range(max(0, cx - radius), min(w, cx + radius + 1)):
                r, g, b, _ = px[x, y]
                corner_colors.append((r, g, b))

    sample_corner(0, 0)
    sample_corner(w - 1, 0)
    sample_corner(0, h - 1)
    sample_corner(w - 1, h - 1)
    bg = Counter(corner_colors).most_common(1)[0][0]

    for y in range(h):
        for x in range(w):
            r, g, b, a = px[x, y]
            if (
                abs(r - bg[0]) <= tolerance
                and abs(g - bg[1]) <= tolerance
                and abs(b - bg[2]) <= tolerance
            ):
                px[x, y] = (r, g, b, 0)

    return img


class MotifCache:
    """Rasterize style-owned SVG assets to PNG for python-pptx placement."""

    def __init__(self, asset_root: Path, out_dir: Path):
        self.asset_root = asset_root
        self.out_dir = out_dir
        self.out_dir.mkdir(parents=True, exist_ok=True)

    def resolve(self, filename: str, max_px: int = 1400, alpha: float = 1.0) -> Path:
        alpha_key = int(round(alpha * 100))
        digest = hashlib.sha256(
            f"{RASTER_CACHE_VERSION}:{filename}:{max_px}:{alpha_key}".encode()
        ).hexdigest()[:16]
        out_path = self.out_dir / f"{digest}.png"
        if out_path.exists():
            return out_path

        svg_path = self.asset_root / filename
        if not svg_path.exists():
            raise FileNotFoundError(f"SVG asset not found: {svg_path}")

        drawing = svg2rlg(str(svg_path))
        if drawing is None:
            raise ValueError(f"Could not parse SVG: {svg_path}")

        scale = min(max_px / max(drawing.width, 1), max_px / max(drawing.height, 1))
        drawing.width *= scale
        drawing.height *= scale
        drawing.scale(scale, scale)

        img = renderPM.drawToPIL(drawing)
        img = strip_solid_background(img)
        if alpha < 0.999:
            r, g, b, a = img.split()
            a = a.point(lambda value: int(value * alpha))
            img.putalpha(a)
        img.save(out_path)
        return out_path


def apply_motifs(slide, motifs: list[dict], cache: MotifCache) -> None:
    for motif in motifs:
        png = cache.resolve(
            motif["file"],
            max_px=motif.get("max_px", 1400),
            alpha=motif.get("alpha", 1.0),
        )
        slide.shapes.add_picture(
            str(png),
            inch(motif["x"]),
            inch(motif["y"]),
            width=inch(motif["w"]),
            height=inch(motif["h"]),
        )


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def inch(value: float):
    return Inches(value)


def resolve_color(color: str) -> str:
    return COLORS.get(color, color)


def set_shape_fill(shape, color: str, transparency: float | int = 0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(resolve_color(color))
    try:
        shape.fill.transparency = transparency
    except Exception:
        pass


def set_shape_line(shape, color: str, width: float = 0.7, transparency: float | int = 0):
    shape.line.color.rgb = rgb(resolve_color(color))
    shape.line.width = Pt(width)
    try:
        shape.line.transparency = transparency
    except Exception:
        pass


def add_text(slide, text, x, y, w, h, size=18, color="ink", bold=False, font="body", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, auto_fit=False, line_spacing=1.08):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = inch(0.03)
    tf.margin_right = inch(0.03)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    if auto_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = FONTS.get(font, font)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(resolve_color(color))
    return box


def add_multiline(slide, lines, x, y, w, h, size=14, color="ink", bullet=True, font="body", leading=1.08, auto_fit=True):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if auto_fit else MSO_AUTO_SIZE.NONE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = leading
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = f"• {line}" if bullet else line
        r.font.name = FONTS.get(font, font)
        r.font.size = Pt(size)
        r.font.color.rgb = rgb(resolve_color(color))
    return box


def add_rect(slide, x, y, w, h, fill="white", line="silver", radius=False, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    set_shape_fill(shape, fill, transparency)
    set_shape_line(shape, line, 0.55, 20)
    return shape


def add_line(slide, x1, y1, x2, y2, color="indigo", width=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    ln.line.color.rgb = rgb(resolve_color(color))
    ln.line.width = Pt(width)
    return ln


def add_chip(slide, x, y, label, color, w=1.16, h=0.31, text_color="white"):
    add_rect(slide, x, y, w, h, fill=color, line=color, radius=True)
    add_text(slide, label, x + 0.08, y + 0.04, w - 0.16, h - 0.04, 7.6, text_color, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_title(slide, title, subtitle=""):
    add_text(slide, title, 0.8, 0.43, 10.7, 0.47, 21, "ink", True, "title")
    add_line(slide, 0.8, 1.02, 3.55, 1.02, "indigo_deep", 1.6)
    add_line(slide, 3.62, 1.02, 4.23, 1.02, "gold", 1.6)
    if subtitle:
        add_text(slide, subtitle, 0.82, 1.08, 9.75, 0.35, 9.5, "ink_soft", False, "body")


def add_footer(slide, idx, total):
    add_line(slide, 0.72, 7.03, 12.63, 7.03, "silver", 0.4)
    add_text(slide, "example: skill_intro_deck / chinese_traditional_color_style", 0.78, 7.09, 6.2, 0.2, 6.7, "ink_soft")
    add_text(slide, f"{idx:02d} / {total:02d}", 11.65, 7.09, 0.9, 0.2, 6.7, "ink_soft", align=PP_ALIGN.RIGHT)


def add_background(slide, texture_path: Path, rail: bool = True):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, inch(W), inch(H))
    set_shape_fill(bg, "paper")
    bg.line.fill.background()
    slide.shapes.add_picture(str(texture_path), 0, 0, width=inch(W), height=inch(H))
    if rail:
        add_rect(slide, 0.28, 0.36, 0.08, 6.75, fill="indigo_deep", line="indigo_deep")
        add_rect(slide, 0.39, 0.36, 0.025, 6.75, fill="gold", line="gold")


def generate_texture() -> Path:
    path = OUT_DIR / "rice-paper-overlay.png"
    if path.exists():
        return path
    random.seed(1337)
    img = Image.new("RGBA", (1920, 1080), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img, "RGBA")
    for _ in range(1200):
        x, y = random.randint(0, 1919), random.randint(0, 1079)
        r = random.choice([1, 1, 2])
        alpha = random.randint(7, 16)
        tone = random.choice([(80, 97, 109), (117, 135, 138), (22, 24, 35)])
        draw.ellipse((x, y, x + r, y + r), fill=(*tone, alpha))
    img = img.filter(ImageFilter.GaussianBlur(0.2))
    img.save(path)
    return path


def add_cover(slide, texture, spec, meta):
    add_background(slide, texture, rail=False)
    add_rect(slide, 0.0, 0.0, 0.28, 7.5, fill="indigo_deep", line="indigo_deep")
    add_text(slide, "my-design-style", 0.86, 0.8, 5.8, 0.55, 14, "indigo_deep", True, "mono")
    add_text(slide, meta["title"], 0.84, 1.42, 7.6, 1.06, 28, "ink", True, "title", line_spacing=0.95)
    add_text(slide, spec.get("tagline", meta["subtitle"]), 0.88, 2.72, 6.85, 0.7, 15, "ink_soft")
    chip_colors = ["indigo_deep", "bamboo", "cang", "rouge"]
    for i, label in enumerate(spec.get("chips", [])):
        add_chip(slide, 0.9 + i * 1.28, 4.08, label, chip_colors[i % len(chip_colors)], 1.02, 0.33)


def add_agenda(slide, texture, spec, idx, total):
    add_background(slide, texture)
    add_title(slide, spec.get("title", "目录"), spec.get("subtitle", ""))
    left, top = 0.95, 1.72
    for i, item in enumerate(spec["items"]):
        head, body = item[0], item[1]
        y = top + i * 0.78
        add_rect(slide, left, y, 0.42, 0.42, fill="indigo_deep" if i < 3 else "cang", line="indigo_deep")
        add_text(slide, str(i + 1), left, y + 0.07, 0.42, 0.18, 9.2, "white", True, "mono", align=PP_ALIGN.CENTER)
        add_text(slide, head, left + 0.62, y - 0.03, 4.0, 0.23, 14.5, "ink", True)
        add_text(slide, body, left + 0.62, y + 0.27, 8.8, 0.22, 9.0, "ink_soft")
    add_footer(slide, idx, total)


def add_section(slide, texture, spec, idx, total):
    add_background(slide, texture, rail=False)
    add_rect(slide, 0.0, 0.0, W, 0.12, fill="indigo_deep", line="indigo_deep")
    add_text(slide, spec["no"], 0.9, 1.04, 1.75, 0.82, 34, "indigo_deep", True, "mono")
    add_text(slide, spec["title"], 0.95, 2.52, 8.9, 0.68, 28, "ink", True, "title")
    add_text(slide, spec["subtitle"], 0.98, 3.45, 7.2, 0.72, 14.2, "ink_soft")
    add_footer(slide, idx, total)


def add_cards(slide, texture, spec, idx, total):
    cols = spec.get("cols", 2)
    cards = spec["cards"]
    add_background(slide, texture)
    add_title(slide, spec["title"], spec.get("subtitle", ""))
    rows = (len(cards) + cols - 1) // cols
    card_w = (11.55 - (cols - 1) * 0.28) / cols
    card_h = min(1.28, (5.02 - (rows - 1) * 0.28) / max(rows, 1))
    for i, card in enumerate(cards):
        c, r = i % cols, i // cols
        x = 0.9 + c * (card_w + 0.28)
        y = 1.74 + r * (card_h + 0.28)
        accent = card.get("accent", "indigo_deep")
        add_rect(slide, x, y, card_w, card_h, fill=card.get("fill", "white"), line="silver", radius=True)
        add_rect(slide, x, y, 0.08, card_h, fill=accent, line=accent)
        add_text(slide, card["head"], x + 0.22, y + 0.12, card_w - 0.36, 0.22, 12.5, accent, True)
        body = card.get("body", [])
        size = card.get("size", 8.8 if len(body) >= 4 else 9.3)
        add_multiline(slide, body, x + 0.22, y + 0.46, card_w - 0.35, card_h - 0.54, size=size, color="ink_soft", leading=1.0)
    add_footer(slide, idx, total)


def add_flow(slide, texture, spec, idx, total):
    steps = spec["steps"]
    add_background(slide, texture)
    add_title(slide, spec["title"], spec.get("subtitle", ""))
    gap, step_w = 0.18, (11.6 - 0.18 * (len(steps) - 1)) / len(steps)
    for i, step in enumerate(steps):
        x = 0.88 + i * (step_w + gap)
        add_rect(slide, x, 2.0, step_w, 2.55, fill="white", line="silver", radius=True)
        add_text(slide, str(i + 1), x + 0.15, 2.22, 0.42, 0.14, 8.3, "white", True, "mono", align=PP_ALIGN.CENTER)
        add_rect(slide, x + 0.15, 2.14, 0.42, 0.42, fill="indigo_deep" if i % 2 == 0 else "cang", line="indigo_deep")
        add_text(slide, step["head"], x + 0.15, 2.7, step_w - 0.3, 0.35, 12, "ink", True, auto_fit=True)
        add_multiline(slide, step["body"], x + 0.16, 3.17, step_w - 0.32, 1.1, 8.2, "ink_soft", bullet=False, leading=1.0)
    add_footer(slide, idx, total)


def add_table(slide, texture, spec, idx, total):
    columns, rows = spec["columns"], spec["rows"]
    col_widths = spec.get("col_widths") or [11.85 / len(columns)] * len(columns)
    add_background(slide, texture)
    add_title(slide, spec["title"], spec.get("subtitle", ""))
    x, y, table_w, row_h = 0.82, 1.62, 11.85, 5.18 / (len(rows) + 1)
    add_rect(slide, x, y, table_w, row_h, fill="indigo_deep", line="indigo_deep")
    cursor = x
    for i, col in enumerate(columns):
        add_text(slide, col, cursor + 0.06, y + 0.11, col_widths[i] - 0.12, row_h - 0.1, 8.2, "white", True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, auto_fit=True)
        cursor += col_widths[i]
    fs = spec.get("font_size", 7.6)
    for r_idx, row in enumerate(rows):
        yy = y + row_h * (r_idx + 1)
        add_rect(slide, x, yy, table_w, row_h, fill="white" if r_idx % 2 == 0 else "mist", line="silver")
        cursor = x
        for c_idx, cell in enumerate(row):
            add_text(slide, cell, cursor + 0.07, yy + 0.06, col_widths[c_idx] - 0.14, row_h - 0.08, fs, "ink", valign=MSO_ANCHOR.MIDDLE, auto_fit=True)
            cursor += col_widths[c_idx]
    add_footer(slide, idx, total)


def add_palette(slide, texture, spec, idx, total):
    add_background(slide, texture)
    add_title(slide, spec["title"], spec.get("subtitle", ""))
    for g, group in enumerate(spec["groups"]):
        y = 1.68 + g * 0.82
        add_text(slide, group["name"], 0.9, y + 0.08, 2.35, 0.24, 10.4, "ink", True, auto_fit=True)
        add_text(slide, group.get("use", ""), 0.9, y + 0.36, 2.35, 0.22, 7.1, "ink_soft", auto_fit=True)
        x = 3.42
        for label, hexv in group["colors"]:
            add_rect(slide, x, y, 0.58, 0.44, fill=hexv, line=hexv)
            add_text(slide, label, x - 0.04, y + 0.48, 0.66, 0.17, 5.7, "ink_soft", align=PP_ALIGN.CENTER, auto_fit=True)
            x += 0.7
    add_footer(slide, idx, total)


def add_decision_tree(slide, texture, spec, idx, total):
    add_background(slide, texture)
    add_title(slide, spec["title"], spec.get("subtitle", ""))
    for i, (cue, style, desc) in enumerate(spec["nodes"]):
        y = 1.9 + i * 1.35
        add_rect(slide, 1.0, y, 3.2, 0.88, fill="white", line="silver", radius=True)
        add_text(slide, cue, 1.18, y + 0.16, 2.85, 0.24, 9.6, "ink", True, auto_fit=True)
        add_line(slide, 4.28, y + 0.44, 5.24, y + 0.44, "gold", 1.2)
        add_rect(slide, 5.35, y, 2.65, 0.88, fill="indigo_deep" if i == 2 else "cang", line="indigo_deep", radius=True)
        add_text(slide, style, 5.48, y + 0.25, 2.35, 0.22, 8.3, "white", True, "mono", align=PP_ALIGN.CENTER, auto_fit=True)
        add_line(slide, 8.1, y + 0.44, 8.82, y + 0.44, "gold", 1.2)
        add_rect(slide, 8.95, y, 3.15, 0.88, fill="mist", line="silver", radius=True)
        add_text(slide, desc, 9.12, y + 0.18, 2.8, 0.45, 8.6, "ink_soft", auto_fit=True)
    add_footer(slide, idx, total)


def add_prompt_slide(slide, texture, spec, idx, total):
    add_background(slide, texture)
    add_title(slide, spec["title"], spec.get("subtitle", ""))
    for i, block in enumerate(spec["prompts"]):
        y = 1.58 + i * 1.55
        accent = block.get("accent", "indigo_deep")
        add_rect(slide, 0.9, y, 11.55, 1.22, fill="white", line="silver", radius=True)
        add_rect(slide, 0.9, y, 0.12, 1.22, fill=accent, line=accent)
        add_text(slide, block["head"], 1.18, y + 0.14, 2.2, 0.2, 9.2, accent, True)
        add_text(slide, block["text"], 3.35, y + 0.13, 8.95, 0.88, 7.9, "ink", auto_fit=True, line_spacing=1.0)
    add_footer(slide, idx, total)


def add_style_lock(slide, texture, spec, idx, total):
    add_background(slide, texture)
    add_title(slide, "本示例的 StyleLock", "previewMode: skip — 默认来自 chinese_traditional_color_style。")
    for i, (k, v, d) in enumerate(spec.get("locks", [])):
        y = 1.58 + i * 0.88
        add_rect(slide, 0.9, y, 3.0, 0.56, fill="mist", line="silver", radius=True)
        add_text(slide, k, 1.05, y + 0.17, 2.5, 0.14, 8.6, "indigo_deep", True, align=PP_ALIGN.CENTER, auto_fit=True)
        add_rect(slide, 4.1, y, 3.0, 0.56, fill="white", line="silver", radius=True)
        add_text(slide, v, 4.25, y + 0.13, 2.65, 0.22, 8.1, "ink", True, align=PP_ALIGN.CENTER, auto_fit=True)
        add_text(slide, d, 7.35, y + 0.08, 4.65, 0.34, 8.1, "ink_soft", auto_fit=True)
    add_footer(slide, idx, total)


def split_motifs(motifs: list[dict]) -> tuple[list[dict], list[dict]]:
    back, front = [], []
    for motif in motifs:
        (back if motif.get("alpha", 1.0) <= 0.4 else front).append(motif)
    return back, front


def render_slide(slide, spec, texture, meta, idx, total, motif_cache: MotifCache):
    motifs = spec.get("motifs", [])
    back_motifs, front_motifs = split_motifs(motifs)
    if back_motifs:
        apply_motifs(slide, back_motifs, motif_cache)

    layout = spec["layout"]
    if layout == "cover":
        add_cover(slide, texture, spec, meta)
    elif layout == "agenda":
        add_agenda(slide, texture, spec, idx, total)
    elif layout == "section":
        add_section(slide, texture, spec, idx, total)
    elif layout == "cards":
        add_cards(slide, texture, spec, idx, total)
    elif layout == "flow":
        add_flow(slide, texture, spec, idx, total)
    elif layout == "table":
        add_table(slide, texture, spec, idx, total)
    elif layout == "palette":
        add_palette(slide, texture, spec, idx, total)
    elif layout == "decision_tree":
        add_decision_tree(slide, texture, spec, idx, total)
    elif layout == "prompt":
        add_prompt_slide(slide, texture, spec, idx, total)
    elif layout == "style_lock":
        add_style_lock(slide, texture, spec, idx, total)
    else:
        raise ValueError(f"Unknown layout: {layout}")

    if front_motifs:
        apply_motifs(slide, front_motifs, motif_cache)


def build():
    outline = json.loads(OUTLINE_PATH.read_text(encoding="utf-8"))
    meta = outline["meta"]
    slides = outline["slides"]
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    texture = generate_texture()
    motif_cache = MotifCache(ASSET_ROOT, OUT_DIR / "svg_raster")

    prs = Presentation()
    prs.slide_width = inch(W)
    prs.slide_height = inch(H)
    blank = prs.slide_layouts[6]
    total = len(slides)
    out_path = EXAMPLE_ROOT / meta.get("outputFilename", "my-design-style-intro.pptx")

    used_svgs: set[str] = set()
    for idx, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        render_slide(slide, spec, texture, meta, idx, total, motif_cache)
        for motif in spec.get("motifs", []):
            used_svgs.add(motif["file"])

    prs.save(out_path)
    print(out_path)
    print(f"slides={total}")
    print(f"svg_motifs={len(used_svgs)}")
    for name in sorted(used_svgs):
        print(f"  - {name}")


if __name__ == "__main__":
    build()
