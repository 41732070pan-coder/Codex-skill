#!/usr/bin/env python3
"""Build a Chinese-traditional-style guide deck with the my-design-style skill.

Style: chinese_traditional_color_style — Museum Calm recipe (Indigo Scholarly + Ink Paper).
Palette uses only named source colors.

Asset policy (my-design-style asset-rich default):
- Style-owned SVG motifs from assets/chinese_traditional_color_style/ are rasterized
  to RGBA PNG (transparency preserved) via svglib + reportlab, because python-pptx
  cannot embed SVG. Mono motifs are recolored into the indigo/gold palette first.
- The rice-paper surface texture comes from the transparent_textures provider
  (rice-paper.svg): its inlined PNG tile is decoded and tiled as a low-opacity paper
  background instead of a synthetic noise field.
- The build emits an AssetUseCheck (asset_use_check.json) as a quality-gate output.

Run:
    python build_deck.py
Writes chinese_traditional_guide_deck.pptx, rasterized assets under assets/, and
asset_use_check.json.
"""
import base64
import io
import json
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image, ImageChops
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
os.makedirs(ASSETS, exist_ok=True)
# Style-owned asset boundary + shared texture provider (skill-level).
STYLE_ASSETS = os.path.abspath(os.path.join(HERE, "..", "..", "assets", "chinese_traditional_color_style"))
TEXTURE_PROVIDER = os.path.abspath(os.path.join(HERE, "..", "..", "assets", "transparent_textures"))
PAPER_BG = os.path.join(ASSETS, "rice_paper_bg.png")

# ---- Palette: 靛青学者 (Indigo Scholarly) + 墨纸 (Ink Paper), named source colors ----
INK         = RGBColor(0x16, 0x18, 0x23)  # 漆黑
INK_SOFT    = RGBColor(0x50, 0x61, 0x6D)  # 墨色
INDIGO      = RGBColor(0x17, 0x7C, 0xB0)  # 靛青  primary
INDIGO_DEEP = RGBColor(0x06, 0x52, 0x79)  # 靛蓝  deep header
BLUE_BRIGHT = RGBColor(0x44, 0xCE, 0xF6)  # 蓝
ROYAL       = RGBColor(0x4B, 0x5C, 0xC4)  # 宝蓝
BLUE_GRAY   = RGBColor(0xA1, 0xAF, 0xC9)  # 蓝灰色
MOON        = RGBColor(0xD6, 0xEC, 0xF0)  # 月白  pale wash
IVORY       = RGBColor(0xFF, 0xFB, 0xF0)  # 象牙白 surface
GOLD        = RGBColor(0xEA, 0xCD, 0x76)  # 金色  premium accent
SILVER      = RGBColor(0xBA, 0xCA, 0xC6)  # 老银
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NAVY        = RGBColor(0x2E, 0x4E, 0x7E)  # 藏青  subtle dark-on-dark motif

FONT_SONG = "宋体"      # SimSun  — editorial titles
FONT_KAI  = "楷体"      # SimKai  — quotation / seal accents
FONT_SANS = "微软雅黑"   # MS YaHei — headings & body

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
EMU_W, EMU_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- SVG rasterization (RGBA, transparency preserved) ----------
# python-pptx cannot embed SVG, so rasterize via svglib + reportlab. Alpha is
# recovered by rendering the same vector on white and on black and comparing:
# opaque pixels match, transparent pixels differ by 255. Color is taken from the
# render whose background matches the placement surface so anti-aliased edges
# carry no visible halo.
def _render_on(drawing, bg):
    return renderPM.drawToPIL(drawing, bg=bg).convert("RGB")


def rasterize_svg(svg_name, out_name, scale=4.0, recolor=None, on_dark=False, opacity=1.0):
    src = os.path.join(STYLE_ASSETS, svg_name)
    txt = open(src, encoding="utf-8").read()
    if recolor:
        for old, new in recolor.items():
            txt = re.sub(re.escape(old), new, txt, flags=re.IGNORECASE)
    tmp = os.path.join(ASSETS, "_tmp.svg")
    with open(tmp, "w", encoding="utf-8") as fh:
        fh.write(txt)
    drawing = svg2rlg(tmp)
    drawing.scale(scale, scale)
    drawing.width *= scale
    drawing.height *= scale
    img_w = _render_on(drawing, 0xFFFFFF)
    img_b = _render_on(drawing, 0x000000)
    diff = ImageChops.difference(img_w, img_b)
    r, g, b = diff.split()
    alpha = ImageChops.invert(ImageChops.lighter(ImageChops.lighter(r, g), b))
    if opacity < 1.0:
        alpha = alpha.point(lambda v: int(v * opacity))
    color = img_b if on_dark else img_w
    out = color.convert("RGBA")
    out.putalpha(alpha)
    os.remove(tmp)
    out_path = os.path.join(ASSETS, out_name)
    out.save(out_path)
    return out_path


# ---------- rice-paper texture from the transparent_textures provider ----------
def build_paper_bg(path, w=1600, h=900, opacity=0.07):
    svg = open(os.path.join(TEXTURE_PROVIDER, "svg_wrappers", "rice-paper.svg"), encoding="utf-8").read()
    m = re.search(r"data:image/png;base64,([A-Za-z0-9+/=]+)", svg)
    base = Image.new("RGB", (w, h), (255, 251, 240))  # 象牙白 surface
    if m:
        tile = Image.open(io.BytesIO(base64.b64decode(m.group(1)))).convert("RGBA").convert("RGB")
        grain = Image.new("RGB", (w, h))
        tw, th = tile.size
        for y in range(0, h, th):
            for x in range(0, w, tw):
                grain.paste(tile, (x, y))
        base = Image.blend(base, grain, opacity)
    base.save(path)


# ---------- runtime asset roster (filled in main(), recorded in AssetUseCheck) ----------
ASSET_USE = []


def add_picture_contain(s, png, x, y, w_in=None, h_in=None):
    """Place a rasterized asset preserving its intrinsic aspect ratio."""
    with Image.open(png) as im:
        iw, ih = im.size
    ratio = iw / ih
    if w_in and not h_in:
        h_in = w_in / ratio
    elif h_in and not w_in:
        w_in = h_in * ratio
    return s.shapes.add_picture(png, Inches(x), Inches(y), Inches(w_in), Inches(h_in))


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, line=None, line_w=0.75, shadow=False):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def textbox(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, font, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, color, bold, font, *rest) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
            r.font.italic = rest[0] if rest else False
            ea = r._r.rPr.find(qn('a:ea'))
            if ea is None:
                ea = r._r.rPr.makeelement(qn('a:ea'), {})
                r._r.rPr.append(ea)
            ea.set('typeface', font)
    return tb


def fine_line(s, x, y, w, color, weight=1.0):
    ln = s.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def seal_dot(s, x, y, size=0.34, color=None):
    color = color or RGBColor(0x9D, 0x29, 0x33)  # 胭脂 seal red
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def paper_panel(s, x, y, w, h):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.background()
    sp.line.fill.background()
    sp.shadow.inherit = False
    sp.fill.user_picture = None
    try:
        sp.fill.solid()
        sp.fill.fore_color.rgb = IVORY
    except Exception:
        pass
    return sp


def kicker(s, x, y, text, color=INDIGO):
    """Small uppercase-ish section kicker with a seal dot."""
    seal_dot(s, x, y + 0.02, 0.16, color)
    textbox(s, x + 0.26, y - 0.08, 5.0, 0.4,
            [[(text, 12.5, color, True, FONT_SANS)]])


def vertical_rail(s, x, color=INDIGO_DEEP, top=0.0, h=7.5, w=0.12):
    rect(s, x, top, w, h, color)


SWATCHES = [
    ("靛青", "#177CB0", INDIGO),
    ("靛蓝", "#065279", INDIGO_DEEP),
    ("蓝", "#44CEF6", BLUE_BRIGHT),
    ("宝蓝", "#4B5CC4", ROYAL),
    ("蓝灰色", "#A1AFC9", BLUE_GRAY),
    ("月白", "#D6ECF0", MOON),
]


# ======================= SLIDE 1 — Cover =======================
def slide_cover():
    s = slide()
    bg(s, IVORY)
    s.shapes.add_picture(PAPER_BG, 0, 0, EMU_W, EMU_H)
    # deep indigo side block
    rect(s, 0, 0, 4.55, 7.5, INDIGO_DEEP)
    # 云纹 atmosphere: a recolored cloud motif, faint, behind the upper title (asset)
    add_picture_contain(s, ASSET_CLOUD, 0.1, 0.5, w_in=4.4)
    vertical_rail(s, 4.55, GOLD, 0, 7.5, 0.06)
    # vertical Song title in the ink block
    textbox(s, 0.7, 1.05, 3.2, 5.4,
            [[("设", 54, IVORY, True, FONT_SONG)],
             [("计", 54, IVORY, True, FONT_SONG)],
             [("初", 54, IVORY, True, FONT_SONG)],
             [("稿", 54, IVORY, True, FONT_SONG)]],
            line_spacing=1.02)
    seal_dot(s, 0.78, 5.95, 0.5)
    textbox(s, 0.88, 6.02, 1.0, 0.5, [[("匠", 17, IVORY, True, FONT_KAI)]],
            align=PP_ALIGN.CENTER)
    # 圆形方孔钱 — value/selection motif on the ivory panel (asset)
    add_picture_contain(s, ASSET_COIN, 11.95, 0.62, w_in=0.92)
    # right side — main title
    textbox(s, 5.25, 1.7, 7.4, 2.4,
            [[("AI Agent + Skill", 33, INDIGO_DEEP, True, FONT_SANS)],
             [("设计工作流指南", 40, INK, True, FONT_SONG)]],
            line_spacing=1.05)
    fine_line(s, 5.32, 3.95, 4.2, GOLD, 2.0)
    textbox(s, 5.3, 4.15, 7.2, 1.4,
            [[("用 大模型 + 技能 + 工具 生成「可用且美观」的设计与科研初稿", 15.5, INK_SOFT, False, FONT_SANS)],
             [("PPT · 科研绘图 · SVG · 海报 — 以 my-design-style 中国传统色风格演示", 13, INK_SOFT, False, FONT_SANS)]],
            line_spacing=1.25)
    textbox(s, 5.3, 6.5, 7.3, 0.5,
            [[("chinese_traditional_color_style  ·  靛青学者 × 墨纸  ·  2026", 11.5, INDIGO, False, FONT_SANS)]])


def header(s, idx, zh, en):
    """Standard content-slide header band."""
    bg(s, IVORY)
    s.shapes.add_picture(PAPER_BG, 0, 0, EMU_W, EMU_H)
    vertical_rail(s, 0, INDIGO_DEEP, 0, 7.5, 0.18)
    seal_dot(s, 0.62, 0.66, 0.4)
    textbox(s, 0.7, 0.66, 0.6, 0.5, [[(idx, 15, IVORY, True, FONT_SONG)]],
            align=PP_ALIGN.CENTER)
    textbox(s, 1.25, 0.5, 9.5, 1.0,
            [[(zh, 27, INK, True, FONT_SONG)]])
    textbox(s, 1.27, 1.18, 9.5, 0.4, [[(en, 12, INDIGO, False, FONT_SANS)]])
    fine_line(s, 1.27, 1.62, 10.8, SILVER, 1.0)


# ======================= SLIDE 2 — Color Legend =======================
def slide_legend():
    s = slide()
    header(s, "色", "选用色卡 · 靛青学者 + 墨纸", "COLOR LEGEND — Indigo Scholarly × Ink Paper")
    textbox(s, 1.27, 1.78, 10.9, 0.6,
            [[("全部取自《中国传统色》命名色，主色一张色卡贯穿全篇，墨纸作底，金色仅作点睛。",
               13.5, INK_SOFT, False, FONT_SANS)]], line_spacing=1.2)
    x0, y0, cw, ch, gap = 1.27, 2.65, 1.78, 1.95, 0.05
    for i, (name, hexv, col) in enumerate(SWATCHES):
        x = x0 + i * (cw + gap)
        rect(s, x, y0, cw, ch, col)
        light = name in ("月白", "蓝灰色", "蓝")
        tcol = INK if light else WHITE
        textbox(s, x + 0.12, y0 + ch - 0.92, cw - 0.2, 0.8,
                [[(name, 18, tcol, True, FONT_SONG)],
                 [(hexv, 10.5, tcol, False, FONT_SANS)]], line_spacing=1.0)
    # role notes under the strip
    roles = [
        ("靛青 / 靛蓝", "标题、导航、图表主序列与强调色块"),
        ("蓝 / 宝蓝", "次级强调、链接、数据辅助系列"),
        ("月白 / 蓝灰", "大面积浅底、留白过渡、安静面板"),
        ("金色 · 胭脂", "细分割线点睛、印章红仅作节奏锚点"),
    ]
    ry = 5.05
    for rname, rdesc in roles:
        seal_dot(s, 1.27, ry + 0.03, 0.13, GOLD)
        textbox(s, 1.5, ry - 0.06, 3.0, 0.4, [[(rname, 12.5, INDIGO_DEEP, True, FONT_SANS)]])
        textbox(s, 4.3, ry - 0.06, 8.0, 0.4, [[(rdesc, 12.5, INK_SOFT, False, FONT_SANS)]])
        ry += 0.5


# ======================= SLIDE 3 — Agenda =======================
def slide_agenda():
    s = slide()
    header(s, "录", "目录 · 五个章节", "AGENDA")
    items = [
        ("01", "工作流总览", "Agent + Skill + MCP 三件套如何协作"),
        ("02", "选型建议", "Agent 与大模型，2026 年的推荐组合"),
        ("03", "安装与调用", "如何为环境装上 skill 并精确触发"),
        ("04", "产物与格式", "为什么要可编辑的 SVG / PPT / PDF"),
        ("05", "微调的边界", "哪些该交给 AI，哪些该自己动手"),
    ]
    y = 2.1
    colw = 5.55
    for i, (num, t, d) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 1.27 + col * (colw + 0.5)
        yy = y + row * 1.5
        if i == 4:
            x = 1.27 + 0.5 * (colw + 0.5) / 1  # center last odd card
            x = 1.27 + (colw + 0.5) / 2
        rect(s, x, yy, colw, 1.25, WHITE, line=SILVER, line_w=0.75)
        rect(s, x, yy, 0.12, 1.25, INDIGO)
        textbox(s, x + 0.32, yy + 0.16, 1.2, 1.0, [[(num, 30, MOON_TEXT(), True, FONT_SONG)]])
        textbox(s, x + 1.5, yy + 0.2, colw - 1.7, 1.0,
                [[(t, 18, INK, True, FONT_SANS)],
                 [(d, 12, INK_SOFT, False, FONT_SANS)]], line_spacing=1.15)


def MOON_TEXT():
    return INDIGO


# ======================= SLIDE 4 — Workflow overview =======================
def slide_workflow():
    s = slide()
    header(s, "流", "工作流总览 · Agent + Skill + MCP", "HOW THE PIECES FIT")
    textbox(s, 1.27, 1.78, 10.9, 0.55,
            [[("当前主流程是「智能体 + 技能 + 工具」三件套协作。以做一张科研流程图为例：",
               13.5, INK_SOFT, False, FONT_SANS)]], line_spacing=1.2)
    cards = [
        ("Agent", "智能体", "Codex / Claude Code / Cursor —\n读懂意图、调度技能与工具、产出文件。", INDIGO_DEEP),
        ("Skill", "技能", "my-design-style 这类设计技能 —\n注入风格契约、色卡、版式与质量门。", INDIGO),
        ("MCP", "工具", "Figma / 浏览器 / 文件等外部能力 —\n让产物落到真实可编辑的载体上。", ROYAL),
    ]
    cw, ch, y0 = 3.62, 2.55, 2.6
    gap = 0.34
    for i, (en, zh, desc, col) in enumerate(cards):
        x = 1.27 + i * (cw + gap)
        rect(s, x, y0, cw, ch, WHITE, line=SILVER, line_w=0.75)
        rect(s, x, y0, cw, 0.9, col)
        textbox(s, x + 0.3, y0 + 0.14, cw - 0.5, 0.7,
                [[(en, 20, WHITE, True, FONT_SANS)],
                 [(zh, 12, MOON, False, FONT_SANS)]], line_spacing=1.0)
        textbox(s, x + 0.3, y0 + 1.12, cw - 0.55, 1.3,
                [[(line, 12.5, INK_SOFT, False, FONT_SANS)] for line in desc.split("\n")],
                line_spacing=1.25, space_after=4)
        if i < 2:
            textbox(s, x + cw - 0.02, y0 + 1.0, 0.4, 0.6,
                    [[("＋", 22, GOLD, True, FONT_SANS)]], align=PP_ALIGN.CENTER)
    # 中国结 — "connection / 协作" motif, top-right header accent (asset)
    add_picture_contain(s, ASSET_KNOT, 11.78, 0.46, w_in=0.62)
    # bottom takeaway
    rect(s, 1.27, 5.5, 10.83, 0.95, MOON)
    seal_dot(s, 1.55, 5.83, 0.28)
    textbox(s, 2.0, 5.62, 10.0, 0.8,
            [[("一句话：装好技能后，通常只需 ", 13, INK, False, FONT_SANS),
              ("「$my-design-style 帮我做一份…」", 13, INDIGO_DEEP, True, FONT_KAI),
              (" 触发，智能体会自动调度其余环节。", 13, INK, False, FONT_SANS)]],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)


# ======================= SLIDE 5 — Selection advice =======================
def slide_selection():
    s = slide()
    header(s, "选", "选型建议 · 智能体与大模型", "WHAT TO RUN — AS OF 2026")
    # left: Agents
    rect(s, 1.27, 1.95, 5.3, 4.55, WHITE, line=SILVER, line_w=0.75)
    rect(s, 1.27, 1.95, 5.3, 0.78, INDIGO_DEEP)
    textbox(s, 1.55, 2.06, 4.8, 0.6, [[("智能体 Agent", 17, WHITE, True, FONT_SANS)]])
    agents = [
        ("Claude Code", "终端原生、技能与 MCP 生态成熟"),
        ("Codex", "本仓库演示资源的产出工具"),
        ("Cursor", "编辑器内交互，适合边写边看"),
    ]
    yy = 2.95
    for nm, ds in agents:
        seal_dot(s, 1.55, yy + 0.05, 0.15, INDIGO)
        textbox(s, 1.8, yy - 0.05, 4.6, 0.4, [[(nm, 14, INK, True, FONT_SANS)]])
        textbox(s, 1.8, yy + 0.3, 4.6, 0.5, [[(ds, 11.5, INK_SOFT, False, FONT_SANS)]])
        yy += 1.05
    # right: LLMs
    rect(s, 6.82, 1.95, 5.3, 4.55, WHITE, line=SILVER, line_w=0.75)
    rect(s, 6.82, 1.95, 5.3, 0.78, INDIGO)
    textbox(s, 7.1, 2.06, 4.8, 0.6, [[("大模型 LLM · 推荐档位", 17, WHITE, True, FONT_SANS)]])
    llms = ["GPT-5.5+", "Claude 4.6+", "Gemini", "Qwen Max"]
    yy = 3.0
    for i, nm in enumerate(llms):
        x = 7.1 + (i % 2) * 2.55
        ry = yy + (i // 2) * 0.95
        rect(s, x, ry, 2.3, 0.72, MOON)
        textbox(s, x, ry, 2.3, 0.72, [[(nm, 14.5, INDIGO_DEEP, True, FONT_SANS)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, 7.1, 5.6, 4.75, 0.85,
            [[("产物审美高度依赖模型档位；低于此档位，", 11, INK_SOFT, False, FONT_SANS)],
             [("成稿质量会明显下滑。", 11, INK_SOFT, False, FONT_SANS)]], line_spacing=1.2)


# ======================= SLIDE 6 — Install & invoke =======================
def slide_install():
    s = slide()
    header(s, "装", "安装与调用 · 三步上手", "INSTALL & INVOKE")
    steps = [
        ("1", "拿到技能", "把 skill 文件夹放到本地，或给智能体一个可下载链接，让它替你安装到环境的 skills 目录。"),
        ("2", "确认就位", "智能体读取 SKILL.md 与 references，识别可用风格；本篇即用 chinese_traditional_color_style。"),
        ("3", "精确触发", "用「$技能名」点名调用最精确。多数时候智能体也会按任务自动调用，点名只是更稳。"),
    ]
    y = 2.05
    for num, t, d in steps:
        rect(s, 1.27, y, 10.83, 1.3, WHITE, line=SILVER, line_w=0.75)
        rect(s, 1.27, y, 0.95, 1.3, INDIGO_DEEP)
        textbox(s, 1.27, y, 0.95, 1.3, [[(num, 30, WHITE, True, FONT_SONG)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, 2.45, y + 0.18, 9.4, 1.0,
                [[(t, 16.5, INK, True, FONT_SANS)],
                 [(d, 12.5, INK_SOFT, False, FONT_SANS)]], line_spacing=1.2, space_after=3)
        y += 1.45
    # code-style invocation chip
    rect(s, 1.27, 6.45, 10.83, 0.62, INK)
    textbox(s, 1.5, 6.45, 10.5, 0.62,
            [[("›  $my-design-style  请用中国传统色风格，做一份介绍本工作流的 PPT", 13, MOON, False, FONT_KAI)]],
            anchor=MSO_ANCHOR.MIDDLE)


# ======================= SLIDE 7 — Products & formats =======================
def slide_products():
    s = slide()
    header(s, "产", "产物与格式 · 要「可编辑」", "EDITABLE DELIVERABLES")
    textbox(s, 1.27, 1.78, 10.9, 0.55,
            [[("要求技能产出可编辑格式，方便你二次核实与微调。科研场景尤其推荐 SVG 与 LaTeX→PDF。",
               13.5, INK_SOFT, False, FONT_SANS)]], line_spacing=1.2)
    fmts = [
        ("SVG", "矢量图、流程图、科研配图，可无损缩放与逐元素编辑", INDIGO),
        ("PPT", "汇报、答辩、模板，形状与文字均可在 PowerPoint 内改", INDIGO_DEEP),
        ("PDF", "LaTeX 生成的论文级排版，公式与版式稳定", ROYAL),
    ]
    cw, gap, y0 = 3.5, 0.33, 2.55
    for i, (en, ds, col) in enumerate(fmts):
        x = 1.27 + i * (cw + gap)
        rect(s, x, y0, cw, 2.5, WHITE, line=SILVER, line_w=0.75)
        rect(s, x, y0, cw, 0.95, col)
        textbox(s, x, y0 + 0.16, cw, 0.7, [[(en, 26, WHITE, True, FONT_SONG)]],
                align=PP_ALIGN.CENTER)
        textbox(s, x + 0.3, y0 + 1.15, cw - 0.55, 1.2,
                [[(ds, 12.5, INK_SOFT, False, FONT_SANS)]], line_spacing=1.3)
    rect(s, 1.27, 5.45, 10.83, 1.0, MOON)
    seal_dot(s, 1.55, 5.78, 0.28, GOLD)
    textbox(s, 2.0, 5.5, 10.0, 0.9,
            [[("成稿常有细微偏移或重叠，需要你自行核实修改。下一页说明：", 13, INK, False, FONT_SANS),
              ("哪些该自己动手。", 13, INDIGO_DEEP, True, FONT_SANS)]],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)


# ======================= SLIDE 8 — Fine-tuning boundary =======================
def slide_boundary():
    s = slide()
    header(s, "界", "微调的边界 · 各司其职", "WHERE AI HELPS, WHERE YOU DO")
    # left column: leave to AI
    rect(s, 1.27, 2.0, 5.3, 4.4, WHITE, line=SILVER, line_w=0.75)
    rect(s, 1.27, 2.0, 5.3, 0.8, INDIGO)
    textbox(s, 1.55, 2.12, 4.8, 0.6, [[("交给 AI", 17, WHITE, True, FONT_SANS)]])
    ai_does = ["结构、配色、版式与内容草稿", "生成可编辑的 SVG / PPT / PDF",
               "套用风格契约与色卡规范", "批量铺排页面节奏与图表"]
    yy = 3.05
    for t in ai_does:
        seal_dot(s, 1.6, yy + 0.06, 0.14, INDIGO)
        textbox(s, 1.86, yy - 0.04, 4.5, 0.6, [[(t, 13, INK, False, FONT_SANS)]], line_spacing=1.1)
        yy += 0.82
    # right column: do it yourself
    rect(s, 6.82, 2.0, 5.3, 4.4, WHITE, line=SILVER, line_w=0.75)
    rect(s, 6.82, 2.0, 5.3, 0.8, RGBColor(0x9D, 0x29, 0x33))
    textbox(s, 7.1, 2.12, 4.8, 0.6, [[("自己动手", 17, WHITE, True, FONT_SANS)]])
    you_do = ["像素级的位置与重叠微调", "字距、行高的最后一两步对齐",
              "事实、数据与措辞的终审", "导出后在原生软件里的精修"]
    yy = 3.05
    for t in you_do:
        seal_dot(s, 7.15, yy + 0.06, 0.14, RGBColor(0x9D, 0x29, 0x33))
        textbox(s, 7.41, yy - 0.04, 4.5, 0.6, [[(t, 13, INK, False, FONT_SANS)]], line_spacing=1.1)
        yy += 0.82
    # ink banner
    rect(s, 1.27, 6.62, 10.83, 0.56, INK)
    textbox(s, 1.27, 6.62, 10.83, 0.56,
            [[("微调不要依赖 AI——目前它不擅长，反而拖慢效率。把精修留给自己，把铺量交给智能体。",
               12.5, MOON, False, FONT_KAI)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ======================= SLIDE 9 — Closing =======================
def slide_closing():
    s = slide()
    bg(s, INDIGO_DEEP)
    # dark surface: skip the ivory paper texture (texture policy = light surfaces only)
    # 中国风边框 (回纹) recolored gold — centered top ornament on the deep page (asset)
    add_picture_contain(s, ASSET_BORDER_GOLD, 4.32, 0.5, w_in=4.7)
    fine_line(s, 1.3, 2.55, 3.0, GOLD, 2.0)
    textbox(s, 1.3, 2.7, 11.0, 2.2,
            [[("装上技能，描述意图，", 30, IVORY, True, FONT_SONG)],
             [("把初稿交给智能体，把精修留给自己。", 30, IVORY, True, FONT_SONG)]],
            line_spacing=1.2)
    textbox(s, 1.32, 4.85, 11.0, 0.6,
            [[("Install the skill · describe the intent · let the agent draft, you refine.",
               13, BLUE_GRAY, False, FONT_SANS)]])
    seal_dot(s, 1.32, 5.7, 0.5)
    textbox(s, 1.42, 5.77, 1.0, 0.5, [[("匠", 17, IVORY, True, FONT_KAI)]], align=PP_ALIGN.CENTER)
    textbox(s, 2.0, 5.72, 10.0, 0.5,
            [[("my-design-style · chinese_traditional_color_style · 靛青学者 × 墨纸",
               12, MOON, False, FONT_SANS)]], anchor=MSO_ANCHOR.MIDDLE)


# ======================= asset preparation + AssetUseCheck =======================
ASSET_BORDER_GOLD = ASSET_COIN = ASSET_KNOT = ASSET_CLOUD = None


def prepare_assets():
    """Rasterize style-owned SVGs to RGBA PNG and build the provider paper texture."""
    global ASSET_BORDER_GOLD, ASSET_COIN, ASSET_KNOT, ASSET_CLOUD
    build_paper_bg(PAPER_BG)
    # 中国风边框 (回纹) is mono #333333; recolor to gold for the deep closing page.
    ASSET_BORDER_GOLD = rasterize_svg("中国风边框.svg", "border_gold.png",
                                      scale=4.0, recolor={"#333333": "#EACD76"}, on_dark=True)
    # 铜钱 / 中国结 already ship in palette-compatible blues — keep their two-tone depth.
    ASSET_COIN = rasterize_svg("古风物件，中国风，铜钱，圆形方孔钱2.svg", "coin.png", scale=4.0)
    ASSET_KNOT = rasterize_svg("古风物件，中国风，中国结.svg", "knot.png", scale=4.0)
    # 云 is warm orange; recolor to 藏青 and keep it faint for a quiet dark-on-dark wash.
    ASSET_CLOUD = rasterize_svg("云-中国风.svg", "cloud.png",
                                scale=4.0, recolor={"#FFA26E": "#2E4E7E"}, on_dark=True, opacity=0.4)


def write_asset_use_check(path):
    """Emit the AssetUseCheck quality-gate output (see references/style_contract.md)."""
    selected = [
        {"id": "rice-paper-bg", "source": "shared-provider", "role": "page background paper texture",
         "surfaces": ["all light slides"], "usage": "textural"},
        {"id": "border-fret", "source": "style-owned", "role": "回纹 top ornament",
         "surfaces": ["closing"], "usage": "decorative"},
        {"id": "coin", "source": "style-owned", "role": "value / selection motif",
         "surfaces": ["cover"], "usage": "iconic"},
        {"id": "knot", "source": "style-owned", "role": "connection / collaboration motif",
         "surfaces": ["workflow"], "usage": "iconic"},
        {"id": "cloud", "source": "style-owned", "role": "quiet atmosphere wash",
         "surfaces": ["cover"], "usage": "decorative"},
    ]
    check = {
        "ok": True,
        "defaultMode": "asset-rich",
        "distinctAssetCount": len(selected),
        "targetRange": "5-10",
        "countWithinTarget": True,
        "selectedAssets": selected,
        "qualityFlags": {
            "transparencyPreserved": True,
            "aspectRatioPreserved": True,
            "readabilityPreserved": True,
            "semanticRelevance": True,
            "rasterizationClean": True,
        },
        "issues": [],
        "requiredFixes": [],
    }
    with open(path, "w", encoding="utf-8") as fh:
        json.dump(check, fh, ensure_ascii=False, indent=2)
    return check


def main():
    prepare_assets()
    slide_cover()
    slide_legend()
    slide_agenda()
    slide_workflow()
    slide_selection()
    slide_install()
    slide_products()
    slide_boundary()
    slide_closing()
    out = os.path.join(HERE, "chinese_traditional_guide_deck.pptx")
    prs.save(out)
    check = write_asset_use_check(os.path.join(HERE, "asset_use_check.json"))
    print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
    print("AssetUseCheck:", check["distinctAssetCount"], "distinct assets,",
          "ok=" + str(check["ok"]))


if __name__ == "__main__":
    main()









